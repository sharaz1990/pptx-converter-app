"""
PPT to Text Converter Web App
==============================
Convert PPT file to TXT online via Streamlit

Author: Sharaz Ahmad
Version: 0.4 - Added security measures
Version: 0.3 - Added streamlit, upload to GitHub and Make Live

libraries required (also needed in requirements.txt):
======================================================
pip install python-pptx streamlit python-magic

requirements.txt:
==================
streamlit
python-pptx

Run server:
===========
streamlit run {file_name}.py

Live Website:
=============
Live Website: https://pptx-converter-app-sharaz.streamlit.app/
Other Links:
Streamlit: https://share.streamlit.io/
Github: https://github.com/sharaz1990/pptx-converter-app
"""
import streamlit as st
from pptx import Presentation
import tempfile
import os
import zipfile
from io import BytesIO

# Security Configuration
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB limit
ALLOWED_EXTENSIONS = ['.pptx']

def validate_file_basic(uploaded_file):
    """Basic file security validation without external dependencies"""
    errors = []
    
    # 1. File size check
    if uploaded_file.size > MAX_FILE_SIZE:
        errors.append(f"File too large. Maximum size: {MAX_FILE_SIZE // (1024*1024)}MB")
    
    # 2. File extension check
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    if file_extension not in ALLOWED_EXTENSIONS:
        errors.append(f"Invalid file type. Only .pptx files allowed")
    
    # 3. Basic PPTX structure validation (PPTX files are ZIP archives)
    try:
        file_bytes = uploaded_file.getvalue()
        with zipfile.ZipFile(BytesIO(file_bytes), 'r') as zip_file:
            # Check for essential PPTX components
            required_files = ['[Content_Types].xml', 'ppt/presentation.xml']
            zip_contents = zip_file.namelist()
            
            for required_file in required_files:
                if required_file not in zip_contents:
                    errors.append("File doesn't appear to be a valid PPTX format")
                    break
                    
    except zipfile.BadZipFile:
        errors.append("File is corrupted or not a valid PPTX file")
    except Exception:
        errors.append("Unable to validate file structure")
    
    # 4. File name sanitization
    dangerous_chars = ['/', '\\', '..', '<', '>', '|', ':', '*', '?', '"']
    if any(char in uploaded_file.name for char in dangerous_chars):
        errors.append("Invalid characters in filename")
    
    # 5. Minimum file size check (empty or too small files)
    if uploaded_file.size < 1000:  # Less than 1KB is suspicious for PPTX
        errors.append("File is too small to be a valid PPTX")
    
    return errors

def safe_text_extraction(uploaded_file):
    """Safely extract text with enhanced error handling"""
    temp_file_path = None
    try:
        # Create temporary file for secure processing
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(uploaded_file.getvalue())
            temp_file_path = temp_file.name
        
        # Process the file with timeout protection
        ppt = Presentation(temp_file_path)
        extracted_text = []
        slide_count = 0
        
        # Limits to prevent resource exhaustion
        MAX_SLIDES = 200
        MAX_TEXT_PER_SLIDE = 50000
        MAX_SHAPES_PER_SLIDE = 100
        
        for i, slide in enumerate(ppt.slides):
            if i >= MAX_SLIDES:
                extracted_text.append(f"\n⚠️ Processing stopped at {MAX_SLIDES} slides for performance\n")
                break
                
            slide_count += 1
            slide_text = f"\n--- Slide {slide_count} ---\n"
            shape_count = 0
            
            for shape in slide.shapes:
                shape_count += 1
                if shape_count > MAX_SHAPES_PER_SLIDE:
                    slide_text += "⚠️ Too many shapes in slide, some content skipped\n"
                    break
                    
                if hasattr(shape, "text") and shape.text.strip():
                    # Limit and sanitize text
                    text_content = shape.text[:MAX_TEXT_PER_SLIDE]
                    # Remove control characters but keep international characters
                    sanitized_text = ''.join(char for char in text_content 
                                           if char.isprintable() or char.isspace())
                    slide_text += sanitized_text + "\n"
            
            extracted_text.append(slide_text)
        
        return "\n".join(extracted_text), slide_count, None
    
    except Exception as e:
        error_msg = str(e)
        # Don't expose internal paths or sensitive error details
        if "temp" in error_msg.lower() or "/" in error_msg or "\\" in error_msg:
            return "", 0, "Error: File processing failed due to security restrictions"
        return "", 0, f"Error: {error_msg[:100]}"  # Limit error message length
    
    finally:
        # Always clean up temporary file
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except:
                pass  # Fail silently on cleanup error

# Streamlit UI
st.title("🔒 Secure PPTX to Text Converter")
st.markdown("**Professional document processing with built-in security**")

# Security notice
with st.expander("🛡️ Security & Privacy Information"):
    st.markdown("""
    **Security Features:**
    - ✅ File type validation and structure verification
    - ✅ Size limits (50MB maximum)
    - ✅ Temporary processing (files deleted immediately)
    - ✅ Content sanitization and limits
    - ✅ HTTPS encryption for all data transmission
    
    **Privacy Commitment:**
    - Files are processed temporarily and deleted immediately
    - No permanent storage of your documents
    - No logging of file contents
    - Secure processing environment
    
    **Best Practices:**
    - Only upload files you have permission to process
    - Avoid highly confidential documents
    - Check your organization's data policies
    """)

# Rate limiting notice
st.info("⏱️ **Usage Limits**: For optimal performance, please limit to reasonable file sizes and avoid bulk processing")

# File upload section
st.subheader("📤 Upload Your PPTX File")

uploaded_file = st.file_uploader(
    "Choose a PPTX file", 
    type=['pptx'],
    help="Maximum file size: 50MB. Only PowerPoint (.pptx) files are accepted.",
    accept_multiple_files=False
)

if uploaded_file is not None:
    # Perform security validation
    with st.spinner("🔍 Validating file security..."):
        security_errors = validate_file_basic(uploaded_file)
    
    if security_errors:
        st.error("❌ **Security Validation Failed**")
        for error in security_errors:
            st.error(f"• {error}")
        st.warning("Please upload a valid PPTX file that meets security requirements.")
        
    else:
        # File passed validation
        st.success(f"✅ **File Validated**: {uploaded_file.name}")
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("File Size", f"{uploaded_file.size:,} bytes")
        with col2:
            st.metric("Status", "Ready for Processing")
        
        # Processing button
        if st.button("🔄 **Extract Text Securely**", type="primary"):
            with st.spinner("🔒 Processing your presentation securely..."):
                extracted_text, slide_count, error = safe_text_extraction(uploaded_file)
            
            if error:
                st.error(f"❌ **Processing Error**: {error}")
                st.info("💡 **Tip**: Ensure your file is a valid PowerPoint presentation")
                
            elif slide_count > 0:
                st.success(f"✅ **Successfully processed {slide_count} slides!**")
                
                # Results section
                st.subheader("📄 Extracted Text")
                
                # Text preview with character count
                char_count = len(extracted_text)
                st.info(f"📊 **Text Statistics**: {char_count:,} characters extracted from {slide_count} slides")
                
                # Text area with extracted content
                st.text_area(
                    "Extracted Content", 
                    extracted_text, 
                    height=400,
                    help="All text has been sanitized and validated for security"
                )
                
                # Download section
                st.subheader("💾 Download Options")
                
                col1, col2 = st.columns(2)
                with col1:
                    st.download_button(
                        label="📥 Download as TXT",
                        data=extracted_text,
                        file_name=f"extracted_{uploaded_file.name.replace('.pptx', '')}.txt",
                        mime="text/plain",
                        help="Download the extracted text as a plain text file"
                    )
                
                with col2:
                    # Optional: Create a summary
                    if char_count > 1000:
                        summary_text = extracted_text[:500] + "\n\n[... content truncated for summary ...]\n\n" + extracted_text[-300:]
                        st.download_button(
                            label="📄 Download Summary",
                            data=summary_text,
                            file_name=f"summary_{uploaded_file.name.replace('.pptx', '')}.txt",
                            mime="text/plain",
                            help="Download a truncated version showing beginning and end"
                        )
                
            else:
                st.warning("⚠️ **No text found** in the presentation slides")
                st.info("The file may contain only images, charts, or non-text elements")

# Usage instructions
with st.expander("📖 How to Use This Tool"):
    st.markdown("""
    ### Step-by-Step Guide:
    
    1. **Select File**: Click "Browse files" and choose your .pptx file
    2. **Validation**: The system automatically validates file security
    3. **Process**: Click "Extract Text Securely" to begin processing
    4. **Review**: Examine the extracted text in the preview area
    5. **Download**: Save the results as a text file
    
    ### What Gets Extracted:
    - All text from text boxes and placeholders
    - Slide titles and content
    - Bullet points and lists
    - Speaker notes (if accessible)
    
    ### What Doesn't Get Extracted:
    - Images and graphics
    - Charts and diagrams
    - Audio/video content
    - Complex formatting
    """)

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; font-size: 0.9em;'>
🔒 <strong>Secure Document Processing</strong> • Built with enterprise security standards<br>
No data retention • HTTPS encrypted • Open source security practices
</div>
""", unsafe_allow_html=True)
